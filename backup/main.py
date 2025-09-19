import os
import uuid
import json
from typing import List, Dict, Any, Tuple
from dataclasses import dataclass

import streamlit as st
from dotenv import load_dotenv

# Document libs
from PyPDF2 import PdfReader
import docx2txt
from pptx import Presentation
import pandas as pd

# Groq client
from groq import Groq
import numpy as np

# --------------------------
# Load env
# --------------------------
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY") or os.getenv("OPENAI_API_KEY") or os.getenv("OPEN_AI_KEY")
if not GROQ_API_KEY:
    st.error("Groq API key not set in environment variables!")
else:
    groq_client = Groq(api_key=GROQ_API_KEY)

# --------------------------
# Embedding model
# --------------------------
EMBEDDING_MODEL = "groq/compound-mini"

# --------------------------
# Helper utilities
# --------------------------
def chunk_text(text: str, max_tokens: int = 500) -> List[str]:
    """Naive text chunker by size (characters). Adjust max_tokens as needed."""
    if not text:
        return []
    approx_chunk_size = max_tokens * 4
    chunks = []
    start = 0
    L = len(text)
    while start < L:
        end = min(start + approx_chunk_size, L)
        if end < L:
            newline = text.rfind('\n', start, end)
            punct = max(text.rfind('.', start, end), text.rfind('?', start, end), text.rfind('!', start, end))
            split_at = max(newline, punct)
            if split_at > start:
                end = split_at + 1
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end
    return chunks

# --------------------------
# MCP message
# --------------------------
@dataclass
class MCPMessage:
    sender: str
    receiver: str
    type: str
    trace_id: str
    payload: Dict[str, Any]

    def to_dict(self):
        return {
            "sender": self.sender,
            "receiver": self.receiver,
            "type": self.type,
            "trace_id": self.trace_id,
            "payload": self.payload,
        }

# --------------------------
# Ingestion agent
# --------------------------
class IngestionAgent:
    """Parses uploaded files into chunks with metadata."""

    def load_file(self, path: str) -> Tuple[List[str], List[Dict[str, Any]]]:
        lower = path.lower()
        try:
            if lower.endswith('.pdf'):
                reader = PdfReader(path)
                text = ""
                for i, page in enumerate(reader.pages):
                    text += f"\n\n[page:{i+1}]\n" + (page.extract_text() or "")
                return self._to_chunks(text, source=os.path.basename(path))

            if lower.endswith('.docx'):
                text = docx2txt.process(path) or ""
                return self._to_chunks(text, source=os.path.basename(path))

            if lower.endswith('.pptx'):
                prs = Presentation(path)
                text = ""
                for si, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text += f"\n\n[slide:{si+1}] " + shape.text
                return self._to_chunks(text, source=os.path.basename(path))

            if lower.endswith('.csv'):
                df = pd.read_csv(path)
                text = df.to_csv(index=False)
                return self._to_chunks(text, source=os.path.basename(path))

            if lower.endswith('.txt') or lower.endswith('.md'):
                with open(path, 'r', encoding='utf-8') as f:
                    text = f.read()
                return self._to_chunks(text, source=os.path.basename(path))

            return [], [{"error": f"Unsupported file type: {path}"}]

        except Exception as e:
            return [], [{"error": str(e)}]

    def _to_chunks(self, text: str, source: str) -> Tuple[List[str], List[Dict[str, Any]]]:
        chunks = chunk_text(text)
        metadata = [{"source": source, "chunk_id": f"{source}::chunk::{i}"} for i in range(len(chunks))]
        return chunks, metadata

# --------------------------
# Retrieval agent (Groq embeddings)
# --------------------------
class RetrievalAgent:
    """Optimized Retrieval Agent using Groq embeddings."""

    def __init__(self):
        self.documents: List[str] = []
        self.metadatas: List[Dict[str, Any]] = []
        self.embeddings: List[np.ndarray] = []

    def add_documents(self, texts: List[str], metadatas: List[Dict[str, Any]]):
        if not texts:
            return {"error": "No texts to add"}
        
        resp = groq_client.embeddings.create(
            model=EMBEDDING_MODEL,
            input=texts
        )
        
        for emb_data in resp['data']:
            emb_vector = np.array(emb_data['embedding'], dtype=np.float32)
            emb_vector /= np.linalg.norm(emb_vector) + 1e-10
            self.embeddings.append(emb_vector)
        
        self.documents.extend(texts)
        self.metadatas.extend(metadatas)
        return {"status": "added", "count": len(texts)}

    def query(self, query_text: str, k: int = 3) -> List[Dict[str, Any]]:
        if not self.embeddings:
            return [{"error": "No documents added yet"}]

        q_emb = groq_client.embeddings.create(
            model=EMBEDDING_MODEL,
            input=query_text
        )['data'][0]['embedding']
        q_vec = np.array(q_emb, dtype=np.float32)
        q_vec /= np.linalg.norm(q_vec) + 1e-10

        sims = [(float(np.dot(q_vec, d_vec)), i) for i, d_vec in enumerate(self.embeddings)]
        sims.sort(key=lambda x: x[0], reverse=True)

        results = []
        for score, idx in sims[:k]:
            results.append({
                "score": score,
                "text": self.documents[idx],
                "metadata": self.metadatas[idx],
            })
        return results

# --------------------------
# LLM agent (Groq)
# --------------------------
class LLMResponseAgent:
    """Calls Groq chat completions with context chunks."""

    def __init__(self, model: str = 'gpt-4o-mini'):
        self.model = model

    def form_prompt(self, query: str, retrieved: List[Dict[str, Any]], chat_history: List[Tuple[str, str]]) -> str:
        parts = ["You are an assistant. Use the context to answer the question."]
        parts.append("\n--- Retrieved chunks ---\n")
        for r in retrieved:
            src = r["metadata"].get("source", "unknown")
            score = r.get("score", 0)
            text = r.get("text", "")
            parts.append(f"[source: {src} | score: {score:.4f}]\n{text}\n")
        parts.append("\n--- Conversation history ---\n")
        for q, a in chat_history[-10:]:
            parts.append(f"User: {q}\nAssistant: {a}\n")
        parts.append(f"\nUser question: {query}\nAnswer:")
        return "\n".join(parts)

    def call_llm(self, prompt: str) -> Dict[str, Any]:
        try:
            resp = groq_client.chat.completions.create(
                model=self.model,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.0,
                max_tokens=512
            )
            text = resp['choices'][0]['message']['content'].strip()
            return {"answer": text, "raw": resp}
        except Exception as e:
            return {"error": str(e)}

    def handle_mcp(self, mcp_msg: MCPMessage) -> MCPMessage:
        payload = mcp_msg.payload
        query = payload.get('query', '')
        retrieved = payload.get('retrieved', [])
        chat_history = payload.get('chat_history', [])
        prompt = self.form_prompt(query, retrieved, chat_history)
        result = self.call_llm(prompt)
        resp_payload = {
            'answer': result.get('answer'),
            'retrieved': retrieved,
            'raw': result.get('raw') if 'raw' in result else None,
            'error': result.get('error') if 'error' in result else None,
        }
        return MCPMessage(sender='LLMResponseAgent', receiver=mcp_msg.sender, type='LLM_ANSWER', trace_id=mcp_msg.trace_id, payload=resp_payload)

# --------------------------
# Streamlit UI
# --------------------------
st.set_page_config(page_title='Agentic RAG Chatbot (Groq)', layout='wide')
st.title('Agentic RAG Chatbot (Groq) for Multi-Format Documents')

with st.sidebar:
    st.header('Upload Documents')
    uploaded_files = st.file_uploader('Upload documents', type=['pdf','docx','pptx','csv','txt','md'], accept_multiple_files=True)
    model_choice = st.selectbox('LLM model', ['gpt-4o-mini', 'gpt-3.5o-mini'])
    k = st.slider('Number of context chunks to retrieve', min_value=1, max_value=8, value=3)

# Initialize agents
if 'ingestion_agent' not in st.session_state:
    st.session_state.ingestion_agent = IngestionAgent()
if 'retrieval_agent' not in st.session_state:
    st.session_state.retrieval_agent = RetrievalAgent()
if 'llm_agent' not in st.session_state:
    st.session_state.llm_agent = LLMResponseAgent(model=model_choice)
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
if 'trace_entries' not in st.session_state:
    st.session_state.trace_entries = []

# Process uploads
if uploaded_files:
    all_texts, all_metas = [], []
    temp_paths = []
    import tempfile
    for file in uploaded_files:
        tmp_dir = tempfile.gettempdir()
        tmp_name = os.path.join(tmp_dir, f"{uuid.uuid4().hex}_{file.name}")
        with open(tmp_name, 'wb') as f:
            f.write(file.getbuffer())
        temp_paths.append(tmp_name)
        texts, metas = st.session_state.ingestion_agent.load_file(tmp_name)
        if texts:
            st.session_state.retrieval_agent.add_documents(texts, metas)
            all_texts.extend(texts)
            all_metas.extend(metas)
        else:
            st.error(f"Failed to parse {file.name}: {metas}")
    if all_texts:
        st.success(f"Ingested {len(all_texts)} chunks from {len(uploaded_files)} files")
    for p in temp_paths:
        try: os.remove(p)
        except: pass

# Chat interface
st.subheader('Ask questions about your documents')
query = st.text_input('Enter your question and press Enter:')

if query:
    retrieved = st.session_state.retrieval_agent.query(query, k=k)
    trace_id = uuid.uuid4().hex
    mcp_in = MCPMessage(sender='Coordinator', receiver='LLMResponseAgent', type='RETRIEVAL_RESULT', trace_id=trace_id, payload={'query': query, 'retrieved': retrieved, 'chat_history': st.session_state.chat_history})
    st.session_state.trace_entries.append(mcp_in.to_dict())

    mcp_out = st.session_state.llm_agent.handle_mcp(mcp_in)
    st.session_state.trace_entries.append(mcp_out.to_dict())

    payload = mcp_out.payload
    if payload.get('error'):
        st.error(f"LLM error: {payload.get('error')}")
    else:
        answer = payload.get('answer', "")
        st.session_state.chat_history.append((query, answer))
        st.markdown('### 💡 Answer')
        st.write(answer)
        st.markdown('---')
        st.markdown('**Source chunks used (top-k)**')
        for r in retrieved:
            src = r['metadata'].get('source', 'unknown')
            st.write(f"- Source: {src} (score: {r.get('score',0):.4f})")
            st.code(r['text'][:1000] + ('...' if len(r['text'])>1000 else ''), language='')

# Show conversation history
if st.session_state.chat_history:
    st.markdown('---')
    st.markdown('### Conversation history')
    for q, a in reversed(st.session_state.chat_history[-10:]):
        st.write(f"**User:** {q}")
        st.write(f"**Assistant:** {a}")
        st.write('---')

# Download MCP trace
if st.sidebar.button('Download MCP trace (.json)'):
    trace_json = json.dumps(st.session_state.trace_entries, indent=2)
    st.sidebar.download_button('Download trace', data=trace_json, file_name='mcp_trace.json', mime='application/json')
