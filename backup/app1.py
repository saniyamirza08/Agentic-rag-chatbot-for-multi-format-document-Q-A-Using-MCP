# app.py
import os
import uuid
import time
import json
from typing import List, Dict, Any
import streamlit as st
from dotenv import load_dotenv

# ✅ Must be the first Streamlit command
st.set_page_config(page_title="RAG with Agents + MCP", layout="wide")

# Load environment variables
load_dotenv()
os.environ['OPEN_API_KEY'] = os.getenv("OPEN_API_KEY", "")
os.environ['GROQ_API_KEY'] = os.getenv("GROQ_API_KEY", "")
os.environ['HF_TOKEN'] = os.getenv("HF_TOKEN", "")


# -------------------------------------------------------------------
# Document Loader Agent
# -------------------------------------------------------------------
class DocumentLoaderAgent:
    def __init__(self):
        pass

    def load_document(self, file_path: str) -> Dict[str, Any]:
        """Loads a document based on its extension."""
        try:
            if file_path.endswith(".pdf"):
                from PyPDF2 import PdfReader
                reader = PdfReader(file_path)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() or ""
                return {"content": text, "type": "pdf"}

            elif file_path.endswith(".docx"):
                import docx2txt
                text = docx2txt.process(file_path)
                return {"content": text, "type": "docx"}

            elif file_path.endswith(".pptx"):
                from pptx import Presentation
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return {"content": text, "type": "pptx"}

            elif file_path.endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as f:
                    return {"content": f.read(), "type": "txt"}

            else:
                return {"error": f"Unsupported file type: {file_path}"}

        except Exception as e:
            return {"error": f"Failed to load {file_path}: {str(e)}"}


# -------------------------------------------------------------------
# Embedding Agent
# -------------------------------------------------------------------
class EmbeddingAgent:
    def __init__(self, model_name="all-MiniLM-L6-v2"):
        from langchain_community.embeddings import HuggingFaceEmbeddings
        self.model_name = model_name
        self.embeddings = HuggingFaceEmbeddings(model_name=self.model_name)

    def embed_text(self, text: str):
        try:
            return self.embeddings.embed_documents([text])[0]
        except Exception as e:
            return {"error": f"Embedding failed: {str(e)}"}


# -------------------------------------------------------------------
# Vector Store Agent
# -------------------------------------------------------------------
class VectorStoreAgent:
    def __init__(self):
        from langchain_community.vectorstores import FAISS
        self.FAISS = FAISS
        self.index = None

    def build_index(self, texts: List[str], embeddings):
        try:
            self.index = self.FAISS.from_texts(texts, embeddings)
            return {"status": "Index built successfully"}
        except Exception as e:
            return {"error": f"Indexing failed: {str(e)}"}

    def query(self, query: str, k=3):
        if not self.index:
            return {"error": "Index not built yet"}
        try:
            results = self.index.similarity_search(query, k=k)
            return {"results": results}
        except Exception as e:
            return {"error": f"Query failed: {str(e)}"}


# -------------------------------------------------------------------
# Query Agent
# -------------------------------------------------------------------
class QueryAgent:
    def __init__(self, llm):
        self.llm = llm

    def answer(self, query: str, context: str):
        try:
            prompt = f"Answer the following question based on the context:\n\nContext:\n{context}\n\nQuestion: {query}\nAnswer:"
            response = self.llm.invoke(prompt)
            return response
        except Exception as e:
            return {"error": f"LLM query failed: {str(e)}"}


# -------------------------------------------------------------------
# Streamlit App
# -------------------------------------------------------------------
def main():
    st.title("Agentic RAG Chatbot for Multi-Format Document QA")

    uploaded_files = st.file_uploader(
        "Upload documents (PDF, DOCX, PPTX, TXT)",
        type=["pdf", "docx", "pptx", "txt"],
        accept_multiple_files=True,
    )

    if uploaded_files:
        doc_loader = DocumentLoaderAgent()
        documents = []
        for file in uploaded_files:
            file_path = f"temp_{uuid.uuid4().hex}_{file.name}"
            with open(file_path, "wb") as f:
                f.write(file.getbuffer())
            result = doc_loader.load_document(file_path)
            if "error" in result:
                st.error(result["error"])
            else:
                documents.append(result["content"])

        if documents:
            st.success(f"✅ Loaded {len(documents)} documents successfully!")

            # Embedding Agent
            embed_agent = EmbeddingAgent()
            st.write("Generating embeddings...")
            embeddings = embed_agent.embeddings  # HuggingFace model

            # Vector Store Agent
            vector_agent = VectorStoreAgent()
            status = vector_agent.build_index(documents, embeddings)
            if "error" in status:
                st.error(status["error"])
                return
            st.success("✅ Index built successfully!")

            # Query Input
            query = st.text_input("Ask a question about your documents:")
            if query:
                results = vector_agent.query(query)
                if "error" in results:
                    st.error(results["error"])
                else:
                    context = "\n".join([r.page_content for r in results["results"]])

                    # Use OpenAI/Groq/Other LLM (dummy example)
                    from langchain_openai import ChatOpenAI
                    llm = ChatOpenAI(model="gpt-3.5-turbo")

                    query_agent = QueryAgent(llm)
                    answer = query_agent.answer(query, context)

                    if isinstance(answer, dict) and "error" in answer:
                        st.error(answer["error"])
                    else:
                        st.subheader("💡 Answer")
                        st.write(answer)
                finally:
        
            if os.path.exists(temp_path):
                os.remove(temp_path)


if __name__ == "__main__":
    main()
