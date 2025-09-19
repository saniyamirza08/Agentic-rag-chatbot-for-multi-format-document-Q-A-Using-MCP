import os
import uuid
import json
from typing import List, Dict, Any, Tuple
from dataclasses import dataclass, field

import streamlit as st
from dotenv import load_dotenv

# Document libs
from PyPDF2 import PdfReader
import docx2txt
from pptx import Presentation
import pandas as pd

# Embeddings & Vector DB
from sentence_transformers import SentenceTransformer
import numpy as np
import faiss

# OpenAI
import openai

# --------------------------
# Load env
# --------------------------
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY") or os.getenv("OPEN_API_KEY") or os.getenv("OPEN_AI_KEY")
if GROQ_API_KEY:
    openai.api_key = GROQ_API_KEY_API_KEY

# --------------------------
# Helper utilities
# --------------------------

def chunk_text(text: str, max_tokens: int = 500) -> List[str]:
    """Naive text chunker by size (characters). Adjust max_tokens as needed."""
    if not text:
        return []
    approx_chunk_size = max_tokens * 4  # heuristic (chars)
    chunks = []
    start = 0
    L = len(text)
    while start < L:
        end = min(start + approx_chunk_size, L)
        # try to break at newline or sentence end
        if end < L:
            newline = text.rfind('\n', start, end)
            punct = max(text.rfind('.', start, end), text.rfind('?', start, end), text.rfind('!', start, end))
            split_at = max(newline, punct)
            if split_at > start:
                end = split_at + 1
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end
    return chunks

# --------------------------
# Agents & MCP
# --------------------------

@dataclass
class MCPMessage:
    sender: str
    receiver: str
    type: str
    trace_id: str
    payload: Dict[str, Any]

    def to_dict(self):
        return {
            "sender": self.sender,
            "receiver": self.receiver,
            "type": self.type,
            "trace_id": self.trace_id,
            "payload": self.payload,
        }


class IngestionAgent:
    """Parses uploaded files and returns a list of document chunks with metadata."""

    def __init__(self):
        pass

    def load_file(self, path: str) -> Tuple[List[str], List[Dict[str, Any]]]:
        lower = path.lower()
        try:
            if lower.endswith('.pdf'):
                reader = PdfReader(path)
                text = ""
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    text += f"\n\n[page:{i+1}]\n" + page_text
                return self._to_chunks(text, source=os.path.basename(path))

            if lower.endswith('.docx'):
                text = docx2txt.process(path) or ""
                return self._to_chunks(text, source=os.path.basename(path))

            if lower.endswith('.pptx'):
                prs = Presentation(path)
                text = ""
                for si, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text += f"\n\n[slide:{si+1}] " + shape.text
                return self._to_chunks(text, source=os.path.basename(path))

            if lower.endswith('.csv'):
                df = pd.read_csv(path)
                text = df.to_csv(index=False)
                return self._to_chunks(text, source=os.path.basename(path))

            if lower.endswith('.txt') or lower.endswith('.md'):
                with open(path, 'r', encoding='utf-8') as f:
                    text = f.read()
                return self._to_chunks(text, source=os.path.basename(path))

            return [], [{"error": f"Unsupported file type: {path}"}]

        except Exception as e:
            return [], [{"error": str(e)}]

    def _to_chunks(self, text: str, source: str) -> Tuple[List[str], List[Dict[str, Any]]]:
        chunks = chunk_text(text)
        metadata = []
        for i, c in enumerate(chunks):
            metadata.append({
                "source": source,
                "chunk_id": f"{source}::chunk::{i}",
                "page": None,
            })
        return chunks, metadata


class RetrievalAgent:
    """Embeds text chunks and builds a FAISS index. Returns top-k chunks on query."""

    def __init__(self, model_name: str = 'all-MiniLM-L6-v2'):
        self.model_name = model_name
        self.embedder = SentenceTransformer(model_name)
        self.index = None
        self.embeddings = None
        self.docs = []
        self.metadatas = []

    def add_documents(self, texts: List[str], metadatas: List[Dict[str, Any]]):
        if not texts:
            return {"error": "No texts to add"}
        vecs = self.embedder.encode(texts, convert_to_numpy=True)
        vecs = vecs.astype('float32')
        if self.index is None:
            d = vecs.shape[1]
            self.index = faiss.IndexFlatIP(d)  # inner product (use normalized vectors)
            # normalize
            faiss.normalize_L2(vecs)
            self.index.add(vecs)
            self.embeddings = vecs
        else:
            faiss.normalize_L2(vecs)
            self.index.add(vecs)
            self.embeddings = np.vstack([self.embeddings, vecs])
        self.docs.extend(texts)
        self.metadatas.extend(metadatas)
        return {"status": "added", "count": len(texts)}

    def query(self, query_text: str, k: int = 3) -> List[Dict[str, Any]]:
        if self.index is None:
            return [{"error": "Index is empty"}]
        qvec = self.embedder.encode([query_text], convert_to_numpy=True).astype('float32')
        faiss.normalize_L2(qvec)
        D, I = self.index.search(qvec, k)
        results = []
        for score, idx in zip(D[0], I[0]):
            if idx < 0 or idx >= len(self.docs):
                continue
            results.append({
                "score": float(score),
                "text": self.docs[idx],
                "metadata": self.metadatas[idx],
            })
        return results


class LLMResponseAgent:
    """Formats MCP messages and calls OpenAI chat completion to get an answer."""

    def __init__(self, model: str = 'gpt-3.5-turbo'):
        self.model = model

    def form_prompt(self, query: str, retrieved: List[Dict[str, Any]], chat_history: List[Tuple[str, str]]) -> str:
        parts = ["You are an assistant. Use the provided context to answer the user's question. If the answer is not present in the context, say you don't know."]
        parts.append("\n--- Retrieved context chunks (with sources) ---\n")
        for r in retrieved:
            src = r["metadata"].get("source", "unknown")
            score = r.get("score", 0)
            text = r.get("text", "")
            parts.append(f"[source: {src} | score: {score:.4f}]\n{text}\n")
        parts.append("\n--- Conversation history ---\n")
        for q, a in chat_history[-10:]:
            parts.append(f"User: {q}\nAssistant: {a}\n")
        parts.append(f"\nUser question: {query}\nAnswer:")
        return "\n".join(parts)

    def call_llm(self, prompt: str, temperature: float = 0.0) -> Dict[str, Any]:
        if not openai.api_key:
            return {"error": "OpenAI API key not set. Set OPENAI_API_KEY in environment."}
        try:
            response = openai.ChatCompletion.create(
                model=self.model,
                messages=[{"role": "user", "content": prompt}],
                temperature=temperature,
                max_tokens=512,
            )
            text = response['choices'][0]['message']['content'].strip()
            return {"answer": text, "raw": response}
        except Exception as e:
            return {"error": str(e)}

    def handle_mcp(self, mcp_msg: MCPMessage) -> MCPMessage:
        # Expect payload to contain: query, retrieved_context (list of chunks)
        payload = mcp_msg.payload
        query = payload.get('query', '')
        retrieved = payload.get('retrieved', [])
        chat_history = payload.get('chat_history', [])
        prompt = self.form_prompt(query, retrieved, chat_history)
        result = self.call_llm(prompt)
        resp_payload = {
            'answer': result.get('answer'),
            'retrieved': retrieved,
            'raw': result.get('raw') if 'raw' in result else None,
            'error': result.get('error') if 'error' in result else None,
        }
        return MCPMessage(sender='LLMResponseAgent', receiver=mcp_msg.sender, type='LLM_ANSWER', trace_id=mcp_msg.trace_id, payload=resp_payload)

# --------------------------
# Streamlit UI
# --------------------------

st.set_page_config(page_title='Agentic RAG Chatbot', layout='wide')
st.title('Agentic RAG Chatbot for Multi-Format Document')

with st.sidebar:
    st.header('Upload Documents')
    uploaded_files = st.file_uploader('Upload documents (pdf, docx, pptx, csv, txt, md)', type=['pdf','docx','pptx','csv','txt','md'], accept_multiple_files=True)
    model_choice = st.selectbox('LLM model', ['gpt-3.5-turbo', 'gpt-4'] if OPENAI_API_KEY else ['gpt-3.5-turbo'])
    k = st.slider('Number of context chunks to retrieve (k)', min_value=1, max_value=8, value=3)
    #st.markdown('---')
    #st.markdown('Environment')
    #st.text('OPENAI_API_KEY set' if openai.api_key else 'OPENAI_API_KEY NOT set')

# Initialize agents in session_state to persist
if 'ingestion_agent' not in st.session_state:
    st.session_state.ingestion_agent = IngestionAgent()
if 'retrieval_agent' not in st.session_state:
    st.session_state.retrieval_agent = RetrievalAgent()
if 'llm_agent' not in st.session_state:
    st.session_state.llm_agent = LLMResponseAgent(model=model_choice)
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []  # list of (question, answer)
if 'trace_entries' not in st.session_state:
    st.session_state.trace_entries = []

# Process uploads
if uploaded_files:
    all_texts = []
    all_metas = []
    temp_paths = []
    for file in uploaded_files:
        # save to temp and load
        import tempfile

    tmp_dir = tempfile.gettempdir()
    tmp_name = os.path.join(tmp_dir, f"{uuid.uuid4().hex}_{file.name}")
    with open(tmp_name, 'wb') as f:
        f.write(file.getbuffer())

        temp_paths.append(tmp_name)
        texts, metas = st.session_state.ingestion_agent.load_file(tmp_name)
        if texts:
            added = st.session_state.retrieval_agent.add_documents(texts, metas)
            all_texts.extend(texts)
            all_metas.extend(metas)
        else:
            st.error(f"Failed to parse {file.name}: {metas}")
    if all_texts:
        st.success(f"Ingested {len(all_texts)} chunks from {len(uploaded_files)} files")
    # cleanup
    for p in temp_paths:
        try:
            os.remove(p)
        except Exception:
            pass

# Chat interface
st.subheader('Ask questions about your documents')
query = st.text_input('Enter your question and press Enter:')

if query:
    # 1. Retrieval agent returns top-k
    retrieved = st.session_state.retrieval_agent.query(query, k=k)

    trace_id = uuid.uuid4().hex
    mcp_in = MCPMessage(sender='Coordinator', receiver='LLMResponseAgent', type='RETRIEVAL_RESULT', trace_id=trace_id, payload={'query': query, 'retrieved': retrieved, 'chat_history': st.session_state.chat_history})
    st.session_state.trace_entries.append(mcp_in.to_dict())

    # 2. LLM agent handles MCP message
    llm_agent = st.session_state.llm_agent
    llm_agent.model = model_choice
    mcp_out = llm_agent.handle_mcp(mcp_in)
    st.session_state.trace_entries.append(mcp_out.to_dict())

    payload = mcp_out.payload
    if payload.get('error'):
        st.error(f"LLM error: {payload.get('error')}")
    else:
        answer = payload.get('answer', "")
        st.session_state.chat_history.append((query, answer))
        st.markdown('### 💡 Answer')
        st.write(answer)

        st.markdown('---')
        st.markdown('**Source chunks used (top-k)**')
        for r in retrieved:
            src = r['metadata'].get('source', 'unknown')
            st.write(f"- Source: {src} (score: {r.get('score',0):.4f})")
            st.code(r['text'][:1000] + ('...' if len(r['text'])>1000 else ''), language='')

        st.markdown('---')
        if st.checkbox('Show MCP trace for this query'):
            st.json([mcp_in.to_dict(), mcp_out.to_dict()])

# Show conversation so far
if st.session_state.chat_history:
    st.markdown('---')
    st.markdown('### Conversation history')
    for i, (q, a) in enumerate(reversed(st.session_state.chat_history[-10:])):
        st.write(f"**User:** {q}")
        st.write(f"**Assistant:** {a}")
        st.write('---')

# Footer / instructions
st.sidebar.markdown('---')
st.sidebar.markdown('**How it works**')
st.sidebar.markdown('- Upload documents.\n- Ask a question.\n- The app ingests, builds embeddings, retrieves top-k chunks, and the LLM answers using those chunks via an MCP message.')

# Save trace to file button
if st.sidebar.button('Download MCP trace (.json)'):
    trace_json = json.dumps(st.session_state.trace_entries, indent=2)
    st.sidebar.download_button('Download trace', data=trace_json, file_name='mcp_trace.json', mime='application/json')

